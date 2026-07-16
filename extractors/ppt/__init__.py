from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from document_engine.models import Document, ImageElement, TableElement
from document_engine.core.registry import Extractor, register_extractor


class PptxExtractor(Extractor):
    format = "pptx"

    def extract(self, file_path: str) -> Document:
        prs = Presentation(file_path)
        p = Path(file_path)

        document = Document(
            title=p.stem,
            file_path=file_path,
            file_format="pptx",
            file_size_bytes=p.stat().st_size,
        )

        slides_text: list[str] = []
        for slide in prs.slides:
            slide_text: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text.append(shape.text)
                if shape.has_table:
                    table = shape.table
                    headers = [cell.text for cell in table.rows[0].cells] if table.rows else []
                    rows = [[cell.text for cell in row.cells] for row in table.rows[1:]]
                    document.tables.append(TableElement(headers=headers, rows=rows))
            slides_text.append("\n".join(slide_text))

        document.text = "\n---\n".join(slides_text)
        document.raw_metadata["slide_count"] = len(prs.slides)

        return document

    def extract_text(self, file_path: str) -> str:
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)
        return "\n".join(texts)


register_extractor(PptxExtractor())
